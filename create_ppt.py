# -*- coding: utf-8 -*-
"""
#   -------------------------------------------------   #
#
#   -------- Description --------
#   1) read.\template.xlsx  as Pandas Dataframe und execute each
#   line in the template.xlsx (each Befehl is one function)
#   2) read template.pptx
#   3) write new_output.pptx with commands from template.xlsx
#   -------- Run --------
#   creat_ppt.py
#
#
#   Created on Fri Sep 22 20:04:13 2017 by   mom
#   -------------------------------------------------   #
"""
import pandas as pd
import os, re, glob
from pptx import Presentation
#from pptx.util import Inches

xls_template = 'template.xlsx'
ppt_template = 'Variantes.pptx'      # or use latestPPtName(),
pptNew = 'Variantes_v01.pptx'

# ---  ersetzt Variablen in der xls_template
#      {{path}}/slide_1.txt  > ./00B001/slide_1.txt
varDict = {'{{path}}': '01B003'
           }

# --------------------------------------------------------------------
#            Definition der Funktion for PPT
# --------------------------------------------------------------------
def newPage(template=0):
    title_slide_layout = prs.slide_layouts[int(template)]
    global actslide
    actslide = prs.slides.add_slide(title_slide_layout)

def tiff(placeholder=1, pfad='', **other):
    assert os.path.isfile(pfad), 'File {} does not exist'.format(pfad)
    placeholder = actslide.placeholders[int(placeholder)]
    placeholder.insert_picture(pfad)

def write(placeholder=1, text='', textFromFile=''):
    if textFromFile:
        with open(textFromFile) as f:
            text = f.read()
    subtitle = actslide.placeholders[int(placeholder)]
    subtitle.text = text

def title(text=''):
    title = actslide.shapes.title
    title.text = text


#               Befehl in der Excel  :  def-function in Python
functionPPT = {'newPage': newPage,
                'tiff': tiff,
                'write': write,
                'title': title}

# --------------------------------------------------------------------
#            Definition der Funktion for Meta
# --------------------------------------------------------------------
def metatiff(placeholder=1, pfad='', **other):
    print('wweee')


#           Befehl in der Excel  :  def-function in Python
functionMeta = {'tiff': metatiff,
                }
# ------------------------------------------------------------

# --- Definition other Functions
def latestPPtName(searchPattern = '*pptx'):
    pptList = glob.glob(searchPattern)

    actPPtName = pptList[-1]     # last item from List. List is sorted by name
    # --- extract number from name. obacht, works only, if there is only one no
    # in the name?  good:   name_123512.pptx, bad: 190702_name_0123.pptx
    actRevNo = re.search(r'(\d+)', actPPtName).group()
    nextRevNo = str(int(actRevNo) + 1)    # increase number by one increament
    nextPPtName = actPPtName.replace(actRevNo, nextRevNo)

    return actPPtName, nextPPtName

def runCommands(xls, functionDict):
    # --- read xls line by line, index: counting number, actBefehl: [Befehl,
    #     Parameter1(...), Parameter2(...)]
    for index, actBefehl in xls.iterrows():
        # --- füge values als ParameterDict zusammen
        #     Parameter: {Befehl: ['Parameter1(...)', 'Parameter2(...)']}
        actParameterList = actBefehl.values[1:]   # alle spalten nach Befehl\
                                                  # array['Template(1)', nan]
        Parameter = {}
        for actParameter in actParameterList:
            try:
                match = re.findall(r'(\w*)\(\ *(.*)\ *\)', actParameter)
                Parameter.update({match[0][0]:match[0][1]})
            except IndexError:
                pass
            except TypeError:
                pass
        # ------------------------------------------------------------

        ParameterReplaced = {}
        # --- replace Variable {{var}} into the Variable  (from varDict)
        for parKey, parValue in Parameter.items():
            for VarPlaceholder, VarValue in varDict.items():
                parKey = parKey.replace(VarPlaceholder, VarValue)
                parValue = parValue.replace(VarPlaceholder, VarValue)
            ParameterReplaced.update({parKey:parValue})

        # --- run command
        actBefehl = actBefehl.Befehl.replace(' ','')   # delete space

        if actBefehl in functionDict:
            functionDict[actBefehl](**ParameterReplaced)
            print ('    {}   {}    {}'.format(index, actBefehl,ParameterReplaced))
        else:
            print ('    {}   {}    {}     not found'.format(index,
                                            actBefehl,ParameterReplaced))


# --------------------------------------------------------------------
#            Öffne PPT
# --------------------------------------------------------------------
ppt_template, pptNew = latestPPtName('Var*pptx')

print (50*'-')
print ('open Presentation:  {}'.format(ppt_template))
print (50*'-' + '\n')

prs = Presentation(ppt_template)

# --- save title in titleList  (currently not used)
# slide_titles = [] # container foe slide titles
# for slide in prs.slides: # iterate over each slide
    # slide_titles.append(slide.shapes.title.text)

# --------------------------------------------------------------------
#            Öffne Excel
# --------------------------------------------------------------------
xls = pd.read_excel(xls_template, index='Befehl')
xls = xls[~xls.Befehl.str.contains(r'^#')]    # drop comment lines

runCommands(xls, functionMeta)     # --- execute Meta Befehle from xls
runCommands(xls, functionPPT)      # --- execute powerpoint Befehle from xls


print ('\n' + 50*'-')
print ('Save New Presentation:  {}'.format(pptNew))
print (50*'-')
prs.save(pptNew)
